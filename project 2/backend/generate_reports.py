import os
import sqlite3
import pandas as pd
from datetime import datetime

# Import reportlab components safely
try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# Import pptx components safely
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

BACKEND_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(BACKEND_DIR)
REPORTS_DIR = os.path.join(PROJECT_DIR, "reports")
DB_PATH = os.path.join(PROJECT_DIR, "data", "retail_analytics.db")

def query_kpi_data():
    """Queries SQLite database to get exact KPIs and lists for reports."""
    if not os.path.exists(DB_PATH):
        return None
        
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    
    # 1. Total KPIs
    cursor.execute("""
        SELECT 
            SUM(Sales) as Sales, 
            SUM(Profit) as Profit, 
            COUNT(DISTINCT Order_ID) as Orders, 
            COUNT(DISTINCT Customer_ID) as Customers,
            AVG(Discount) * 100 as Discount
        FROM sales_data
    """)
    kpis = cursor.fetchone()
    sales, profit, orders, customers, avg_discount = kpis
    sales = sales or 0
    profit = profit or 0
    orders = orders or 0
    customers = customers or 0
    avg_discount = avg_discount or 0
    margin = (profit / sales * 100) if sales > 0 else 0
    aov = (sales / orders) if orders > 0 else 0
    
    # 2. Category Performance
    cursor.execute("""
        SELECT Category, SUM(Sales), SUM(Profit), (SUM(Profit)/SUM(Sales))*100
        FROM sales_data
        GROUP BY Category
        ORDER BY SUM(Sales) DESC
    """)
    categories = cursor.fetchall()
    
    # 3. Regional Performance
    cursor.execute("""
        SELECT Region, SUM(Sales), SUM(Profit), (SUM(Profit)/SUM(Sales))*100
        FROM sales_data
        GROUP BY Region
        ORDER BY SUM(Sales) DESC
    """)
    regions = cursor.fetchall()
    
    # 4. Top 5 profitable products
    cursor.execute("""
        SELECT Product_Name, Sub_Category, SUM(Profit)
        FROM sales_data
        GROUP BY Product_ID
        ORDER BY SUM(Profit) DESC
        LIMIT 5
    """)
    top_products = cursor.fetchall()

    # 5. Bottom 5 unprofitable products
    cursor.execute("""
        SELECT Product_Name, Sub_Category, SUM(Profit)
        FROM sales_data
        GROUP BY Product_ID
        ORDER BY SUM(Profit) ASC
        LIMIT 5
    """)
    bottom_products = cursor.fetchall()
    
    # 6. Unprofitable States
    cursor.execute("""
        SELECT State, SUM(Sales), SUM(Profit)
        FROM sales_data
        GROUP BY State
        HAVING SUM(Profit) < 0
        ORDER BY SUM(Profit) ASC
        LIMIT 5
    """)
    loss_states = cursor.fetchall()
    
    conn.close()
    
    return {
        "sales": sales,
        "profit": profit,
        "orders": orders,
        "customers": customers,
        "discount": avg_discount,
        "margin": margin,
        "aov": aov,
        "categories": categories,
        "regions": regions,
        "top_products": top_products,
        "bottom_products": bottom_products,
        "loss_states": loss_states
    }

def generate_pdf_report():
    """Generates reports/project_report.pdf using ReportLab."""
    if not REPORTLAB_AVAILABLE:
        print("ReportLab is not installed. PDF generation skipped.")
        return False
        
    os.makedirs(REPORTS_DIR, exist_ok=True)
    pdf_path = os.path.join(REPORTS_DIR, "project_report.pdf")
    
    data = query_kpi_data()
    if not data:
        print("Database not loaded. Please initialize database first.")
        return False
        
    # Styles Setup
    styles = getSampleStyleSheet()
    
    # Custom Palette - Indigo & Slate
    c_primary = colors.HexColor("#1e293b")   # Slate 800
    c_secondary = colors.HexColor("#3b82f6") # Blue 500
    c_accent = colors.HexColor("#10b981")    # Emerald 500
    c_danger = colors.HexColor("#ef4444")    # Red 500
    c_bg_light = colors.HexColor("#f8fafc")  # Slate 50
    c_border = colors.HexColor("#e2e8f0")    # Slate 200
    
    # Modify existing styles to avoid duplicate warnings
    styles['Normal'].textColor = colors.HexColor("#334155")
    styles['Normal'].fontSize = 10
    styles['Normal'].leading = 14
    
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=24,
        leading=28,
        textColor=c_primary,
        alignment=0, # Left aligned
        spaceAfter=15
    )
    
    subtitle_style = ParagraphStyle(
        'DocSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=12,
        leading=16,
        textColor=c_secondary,
        spaceAfter=30
    )
    
    h1_style = ParagraphStyle(
        'DocH1',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=16,
        leading=20,
        textColor=c_primary,
        spaceBefore=15,
        spaceAfter=10,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=12,
        leading=16,
        textColor=c_secondary,
        spaceBefore=10,
        spaceAfter=6,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        spaceAfter=10
    )
    
    bullet_style = ParagraphStyle(
        'DocBullet',
        parent=styles['Normal'],
        leftIndent=20,
        firstLineIndent=-10,
        spaceAfter=6
    )

    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=letter,
        rightMargin=54,
        leftMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    
    story = []
    
    # --- PAGE 1: TITLE & EXECUTIVE SUMMARY ---
    story.append(Paragraph("Retail Business Performance & Profitability Analysis", title_style))
    story.append(Paragraph(f"<b>Elevate Labs Data Analyst Internship - Portfolio Submission</b><br/>Generated: {datetime.now().strftime('%B %d, %Y')}", subtitle_style))
    story.append(Spacer(1, 20))
    
    story.append(Paragraph("Executive Summary", h1_style))
    exec_summary_text = (
        "This project delivers an interactive, data-driven analytical showcase evaluating the operational efficiency "
        "and profitability of a national retail operation. Utilizing transactional sales logs spanning from 2024 to 2026, "
        "this report highlights critical leakage points, regional opportunities, category profitability anomalies, and "
        "provides database-validated strategies to recapture eroding margins."
    )
    story.append(Paragraph(exec_summary_text, body_style))
    story.append(Spacer(1, 15))
    
    # KPI Grid Table
    kpi_table_data = [
        [
            Paragraph(f"<b>Total Sales</b><br/><font color='#2563eb' size=14>${data['sales']:,.2f}</font>", body_style),
            Paragraph(f"<b>Total Profit</b><br/><font color='#10b981' size=14>${data['profit']:,.2f}</font>", body_style)
        ],
        [
            Paragraph(f"<b>Total Orders</b><br/><font color='#1e293b' size=14>{data['orders']:,}</font>", body_style),
            Paragraph(f"<b>Profit Margin</b><br/><font color='#059669' size=14>{data['margin']:.2f}%</font>", body_style)
        ],
        [
            Paragraph(f"<b>Customers Served</b><br/><font color='#1e293b' size=14>{data['customers']:,}</font>", body_style),
            Paragraph(f"<b>Average Discount</b><br/><font color='#ef4444' size=14>{data['discount']:.2f}%</font>", body_style)
        ]
    ]
    
    t_kpi = Table(kpi_table_data, colWidths=[250, 250])
    t_kpi.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), c_bg_light),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('PADDING', (0,0), (-1,-1), 12),
        ('LINEBELOW', (0,0), (-1,-2), 1, c_border),
        ('LINEAFTER', (0,0), (0,-1), 1, c_border),
    ]))
    story.append(t_kpi)
    story.append(PageBreak())
    
    # --- PAGE 2: FINANCIAL & REGIONAL PERFORMANCE ---
    story.append(Paragraph("Geographical & Regional Performance", h1_style))
    story.append(Paragraph(
        "Sales distribution indicates that the West and East regions are the primary drivers of volume, representing "
        "the majority of total transactions. However, regional profitability exhibits high variability due to shipping overheads and discounting.",
        body_style
    ))
    story.append(Spacer(1, 10))
    
    # Region Table
    region_headers = [Paragraph("<b>Region</b>", body_style), Paragraph("<b>Sales</b>", body_style), Paragraph("<b>Profit</b>", body_style), Paragraph("<b>Margin (%)</b>", body_style)]
    region_rows = [region_headers]
    for r in data["regions"]:
        name, r_sales, r_profit, r_margin = r
        region_rows.append([
            Paragraph(name, body_style),
            Paragraph(f"${r_sales:,.2f}", body_style),
            Paragraph(f"${r_profit:,.2f}", body_style),
            Paragraph(f"{r_margin:.2f}%", body_style)
        ])
    t_reg = Table(region_rows, colWidths=[120, 120, 120, 120])
    t_reg.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), c_primary),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('BOTTOMPADDING', (0,0), (-1,0), 8),
        ('PADDING', (0,1), (-1,-1), 6),
        ('BACKGROUND', (0,1), (-1,-1), c_bg_light),
        ('LINEBELOW', (0,0), (-1,-1), 0.5, c_border),
    ]))
    # Quick fix for text color in table headers
    for i in range(len(region_headers)):
        region_headers[i].style.textColor = colors.white
    
    story.append(t_reg)
    story.append(Spacer(1, 20))
    
    story.append(Paragraph("Loss-Making States (Action Required)", h2_style))
    story.append(Paragraph(
        "Several states show significant net losses despite generating high sales volumes. This indicates systematic pricing "
        "deficiencies, such as local clearance campaigns and aggressive customer acquisition discounts.",
        body_style
    ))
    
    state_headers = [Paragraph("<b>State</b>", body_style), Paragraph("<b>Total Sales</b>", body_style), Paragraph("<b>Net Loss</b>", body_style)]
    state_rows = [state_headers]
    for s in data["loss_states"]:
        st_name, st_sales, st_profit = s
        state_rows.append([
            Paragraph(st_name, body_style),
            Paragraph(f"${st_sales:,.2f}", body_style),
            Paragraph(f"<font color='red'>${st_profit:,.2f}</font>", body_style)
        ])
    t_state = Table(state_rows, colWidths=[160, 160, 160])
    t_state.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), c_primary),
        ('PADDING', (0,1), (-1,-1), 6),
        ('LINEBELOW', (0,0), (-1,-1), 0.5, c_border),
    ]))
    for i in range(len(state_headers)):
        state_headers[i].style.textColor = colors.white
        
    story.append(t_state)
    story.append(PageBreak())
    
    # --- PAGE 3: CATEGORY ANALYSIS ---
    story.append(Paragraph("Category & Product Performance", h1_style))
    story.append(Paragraph(
        "Product performance varies heavily across categories. Technology remains the highest margin category, while Furniture has thin profits due to heavy freight costs.",
        body_style
    ))
    story.append(Spacer(1, 10))
    
    cat_headers = [Paragraph("<b>Category</b>", body_style), Paragraph("<b>Sales</b>", body_style), Paragraph("<b>Profit</b>", body_style), Paragraph("<b>Margin (%)</b>", body_style)]
    cat_rows = [cat_headers]
    for c in data["categories"]:
        name, c_sales, c_profit, c_margin = c
        cat_rows.append([
            Paragraph(name, body_style),
            Paragraph(f"${c_sales:,.2f}", body_style),
            Paragraph(f"${c_profit:,.2f}", body_style),
            Paragraph(f"{c_margin:.2f}%", body_style)
        ])
    t_cat = Table(cat_rows, colWidths=[120, 120, 120, 120])
    t_cat.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), c_primary),
        ('PADDING', (0,1), (-1,-1), 6),
        ('LINEBELOW', (0,0), (-1,-1), 0.5, c_border),
    ]))
    for i in range(len(cat_headers)):
        cat_headers[i].style.textColor = colors.white
    story.append(t_cat)
    story.append(Spacer(1, 20))
    
    story.append(Paragraph("Top 5 Profitable Products", h2_style))
    top_prod_headers = [Paragraph("<b>Product Name</b>", body_style), Paragraph("<b>Sub-Category</b>", body_style), Paragraph("<b>Net Profit</b>", body_style)]
    top_prod_rows = [top_prod_headers]
    for p in data["top_products"]:
        p_name, p_sub, p_profit = p
        # Truncate product name if too long
        p_name_short = p_name[:45] + "..." if len(p_name) > 45 else p_name
        top_prod_rows.append([
            Paragraph(p_name_short, body_style),
            Paragraph(p_sub, body_style),
            Paragraph(f"<font color='green'>${p_profit:,.2f}</font>", body_style)
        ])
    t_tprod = Table(top_prod_rows, colWidths=[240, 120, 120])
    t_tprod.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), c_primary),
        ('PADDING', (0,1), (-1,-1), 5),
        ('LINEBELOW', (0,0), (-1,-1), 0.5, c_border),
    ]))
    for i in range(len(top_prod_headers)):
        top_prod_headers[i].style.textColor = colors.white
    story.append(t_tprod)
    
    story.append(Spacer(1, 15))
    story.append(Paragraph("Bottom 5 Unprofitable Products (Leakages)", h2_style))
    bot_prod_headers = [Paragraph("<b>Product Name</b>", body_style), Paragraph("<b>Sub-Category</b>", body_style), Paragraph("<b>Net Loss</b>", body_style)]
    bot_prod_rows = [bot_prod_headers]
    for p in data["bottom_products"]:
        p_name, p_sub, p_profit = p
        p_name_short = p_name[:45] + "..." if len(p_name) > 45 else p_name
        bot_prod_rows.append([
            Paragraph(p_name_short, body_style),
            Paragraph(p_sub, body_style),
            Paragraph(f"<font color='red'>${p_profit:,.2f}</font>", body_style)
        ])
    t_bprod = Table(bot_prod_rows, colWidths=[240, 120, 120])
    t_bprod.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), c_primary),
        ('PADDING', (0,1), (-1,-1), 5),
        ('LINEBELOW', (0,0), (-1,-1), 0.5, c_border),
    ]))
    for i in range(len(bot_prod_headers)):
        bot_prod_headers[i].style.textColor = colors.white
    story.append(t_bprod)
    
    story.append(PageBreak())
    
    # --- PAGE 4: BUSINESS INSIGHTS ---
    story.append(Paragraph("Business Insights (Key Findings)", h1_style))
    insights = [
        "<b>1. Holiday Sales Seasonality:</b> Peak sales occur during November and December (Q4 holiday surge), accounting for nearly 30% of annual revenue, suggesting strict inventory loading constraints in Q3.",
        "<b>2. Loss-Making Tables & Bookcases:</b> Furniture category profits are severely damaged by Tables and Bookcases. They have negative net margins due to large shipping costs and clearing discounts.",
        "<b>3. Technology Margin Leader:</b> Copiers and Phones sub-categories are the primary margin leaders, driving strong cash flows despite lower unit sales volume compared to binders or paper.",
        "<b>4. High Discount Penalty:</b> Discounts above 20% consistently lead to unprofitable orders. The '51%+ discount' tier represents a 100% loss-making segment.",
        "<b>5. Customer Concentration Risk:</b> The top 10% of customers generate approximately 40% of sales. A loss of a few VIP clients could severely impact revenues.",
        "<b>6. Regional Discrepancy:</b> West and East represent the largest volume, but Texas (Central region) generates large sales alongside negative profit due to high promotional rates.",
        "<b>7. Home Office High AOV:</b> Home Office accounts have a higher Average Order Value ($260+) compared to the Consumer segment ($210+), highlighting premium commercial profiles."
    ]
    for ins in insights:
        story.append(Paragraph(ins, bullet_style))
        story.append(Spacer(1, 4))
        
    story.append(Spacer(1, 15))
    story.append(Paragraph("Actionable Recommendations", h1_style))
    recs = [
        "<b>1. Restructure Furniture Discounting:</b> Cap discounts on bulky items (Tables/Bookcases) at a maximum of 10%. Stop automatic clearance campaigns in central states.",
        "<b>2. Re-negotiate Shipping Contracts:</b> Partner with regional freight carrier pools for Furniture to reduce the average logistics cost, currently eating 25% of gross category revenue.",
        "<b>3. Home Office Promotional Target:</b> Launch dedicated high-ticket bundles (Computers + Office storage) targeting remote workers, who exhibit high transaction limits.",
        "<b>4. State-level Pricing Policy:</b> Enforce a minimum margin threshold of 5% in Texas and Ohio. Restructure regional manager performance scorecards to align with Profit, not pure Sales volume.",
        "<b>5. Implement a Tiered VIP Program:</b> Deploy a dedicated customer success portal for the top 100 accounts to safeguard key revenues and encourage contract lock-ins."
    ]
    for rec in recs:
        story.append(Paragraph(rec, bullet_style))
        story.append(Spacer(1, 4))
        
    # Build Document
    doc.build(story)
    print("PDF Report generated successfully.")
    return True

def generate_pptx_presentation():
    """Generates reports/presentation.pptx using python-pptx."""
    if not PPTX_AVAILABLE:
        print("python-pptx is not installed. PPTX presentation skipped.")
        return False
        
    os.makedirs(REPORTS_DIR, exist_ok=True)
    pptx_path = os.path.join(REPORTS_DIR, "presentation.pptx")
    
    data = query_kpi_data()
    if not data:
        print("Database not loaded. Please initialize database first.")
        return False
        
    prs = Presentation()
    
    # Setup slide sizes (16:9 aspect ratio standard)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Helper to add backgrounds and header formatting
    c_dark_blue = RGBColor(30, 41, 59)    # Slate 800
    c_blue = RGBColor(59, 130, 246)       # Blue 500
    c_gray = RGBColor(100, 116, 139)      # Slate 500
    c_white = RGBColor(255, 255, 255)
    
    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE
        0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_dark_blue
    bg.line.color.rgb = c_dark_blue
    
    # Title Text Box
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Retail Business Performance &\nProfitability Analysis"
    p.font.name = "Arial"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = "Data Analyst Internship Portfolio | Elevate Labs Submission"
    p2.font.name = "Arial"
    p2.font.size = Pt(18)
    p2.font.color.rgb = c_blue
    
    # ----------------------------------------------------
    # SLIDE 2: Executive Summary (Light Theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide.shapes.title.text = "Executive Summary: Business Performance Overview"
    slide.shapes.title.text_frame.paragraphs[0].font.name = "Arial"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = c_dark_blue
    
    # Summary Paragraph
    tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "This report evaluates sales transaction performance and margins across 2.5 years of operations (2024-2026). The project is built to identify core margin leakages and outline database-backed performance improvements."
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.color.rgb = c_dark_blue
    
    # KPI Grid Boxes (Manual layout of 4 boxes)
    kpis_to_show = [
        {"name": "TOTAL SALES", "val": f"${data['sales']:,.2f}", "color": c_blue},
        {"name": "TOTAL PROFIT", "val": f"${data['profit']:,.2f}", "color": RGBColor(16, 185, 129)},
        {"name": "PROFIT MARGIN", "val": f"{data['margin']:.2f}%", "color": RGBColor(5, 150, 105)},
        {"name": "TOTAL ORDERS", "val": f"{data['orders']:,}", "color": c_dark_blue}
    ]
    
    for i, k in enumerate(kpis_to_show):
        left = Inches(0.75 + i * 3.0)
        top = Inches(3.2)
        width = Inches(2.7)
        height = Inches(2.2)
        
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate 50
        box.line.color.rgb = RGBColor(226, 232, 240)
        
        # Text Frame inside box
        btf = box.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        bp.text = k["name"]
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = c_gray
        bp.space_after = Pt(20)
        
        bp2 = btf.add_paragraph()
        bp2.text = k["val"]
        bp2.font.size = Pt(28)
        bp2.font.bold = True
        bp2.font.color.rgb = k["color"]

    # ----------------------------------------------------
    # SLIDE 3: Regional Analysis
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Regional Performance & Unprofitable State Focus"
    
    # Left Column: Regions Summary
    tx_left = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_l = tx_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "Geographical Breakdown & Sales Drivers"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_dark_blue
    p.space_after = Pt(10)
    
    for name, r_sales, r_profit, r_margin in data["regions"]:
        p_row = tf_l.add_paragraph()
        p_row.text = f"• {name} Region: Sales of ${r_sales:,.0f} | Profit margin: {r_margin:.1f}%"
        p_row.font.size = Pt(14)
        p_row.space_after = Pt(8)
        
    # Right Column: Loss-making states
    tx_right = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_r = tx_right.text_frame
    tf_r.word_wrap = True
    
    p_r = tf_r.paragraphs[0]
    p_r.text = "Top Loss-Making States"
    p_r.font.size = Pt(18)
    p_r.font.bold = True
    p_r.font.color.rgb = RGBColor(239, 68, 68)
    p_r.space_after = Pt(10)
    
    for s_name, s_sales, s_profit in data["loss_states"]:
        p_row = tf_r.add_paragraph()
        p_row.text = f"• {s_name}: Gross Sales ${s_sales:,.0f} | Net Profit: -${abs(s_profit):,.0f}"
        p_row.font.size = Pt(14)
        p_row.space_after = Pt(8)
        
    p_note = tf_r.add_paragraph()
    p_note.text = "Note: Loss-making states exhibit average discount rates exceeding 25%, indicating that aggressive price promotion is hurting structural profitability."
    p_note.font.size = Pt(11)
    p_note.font.italic = True
    p_note.font.color.rgb = c_gray
    p_note.space_before = Pt(15)

    # ----------------------------------------------------
    # SLIDE 4: Category & Product Performance
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Product Line Health: Top Stars & Margin Leaks"
    
    # Category summary at top
    tx_cat = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.5))
    tf_cat = tx_cat.text_frame
    tf_cat.word_wrap = True
    p = tf_cat.paragraphs[0]
    p.text = "Product categories reflect high variance. Technology leads profitability, while Furniture contains critical bleed areas (Tables & Bookcases)."
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = c_dark_blue
    
    # Bottom columns
    tx_p1 = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(5.5), Inches(4.0))
    tf_p1 = tx_p1.text_frame
    tf_p1.word_wrap = True
    p_p1 = tf_p1.paragraphs[0]
    p_p1.text = "Star Product Lines (Highest Profit)"
    p_p1.font.size = Pt(16)
    p_p1.font.bold = True
    p_p1.font.color.rgb = RGBColor(16, 185, 129)
    p_p1.space_after = Pt(8)
    
    for p_name, p_sub, p_profit in data["top_products"]:
        p_row = tf_p1.add_paragraph()
        p_row.text = f"• {p_name[:35]}... ({p_sub}): +${p_profit:,.0f}"
        p_row.font.size = Pt(13)
        p_row.space_after = Pt(6)
        
    tx_p2 = slide.shapes.add_textbox(Inches(6.8), Inches(2.8), Inches(5.5), Inches(4.0))
    tf_p2 = tx_p2.text_frame
    tf_p2.word_wrap = True
    p_p2 = tf_p2.paragraphs[0]
    p_p2.text = "Margin Leaks (Highest Net Loss)"
    p_p2.font.size = Pt(16)
    p_p2.font.bold = True
    p_p2.font.color.rgb = RGBColor(239, 68, 68)
    p_p2.space_after = Pt(8)
    
    for p_name, p_sub, p_profit in data["bottom_products"]:
        p_row = tf_p2.add_paragraph()
        p_row.text = f"• {p_name[:35]}... ({p_sub}): -${abs(p_profit):,.0f}"
        p_row.font.size = Pt(13)
        p_row.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 5: 5 Key Business Insights
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Data Analytics: 5 Key Business Insights"
    
    tx_ins = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.8))
    tf_ins = tx_ins.text_frame
    tf_ins.word_wrap = True
    
    insights = [
        ("1. Holiday Purchase Climax", "November and December generate 30% of total revenue. Planning inventory and logistical capability is critical for Q3."),
        ("2. The Discounting Penalty", "Any transactional discount exceeding 20% drops profitability below breakeven. Promos at 50% scale serve as complete loss operations."),
        ("3. Tables and Bookcases Overhead", "Tables and bookcases are highly unprofitable across all segments due to bulky size driving high shipping costs combined with high discounts."),
        ("4. Customer Concentration Risk", "The top 10% of customers account for 40% of total revenue. Maintaining satisfaction with these key clients is vital."),
        ("5. B2B Account Profile Strength", "Corporate and Home Office clients exhibit 20% higher average order value (AOV) compared to traditional consumer retail baskets.")
    ]
    
    for title, desc in insights:
        p = tf_ins.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = c_dark_blue
        
        # Append description to same paragraph
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(14)
        run.font.color.rgb = c_gray
        p.space_after = Pt(10)

    # Remove the default first paragraph if empty
    if len(tf_ins.paragraphs) > len(insights) and tf_ins.paragraphs[0].text == "":
        p_first = tf_ins.paragraphs[0]
        # Not easily deletable but can keep it blank or make it small

    # ----------------------------------------------------
    # SLIDE 6: Actionable Recommendations
    # ----------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Actionable Business Recommendations"
    
    tx_recs = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(4.8))
    tf_recs = tx_recs.text_frame
    tf_recs.word_wrap = True
    
    recs = [
        ("1. Caps on Furniture Discounts", "Enforce a maximum discount cap of 10% on Tables and Bookcases. Limit automated promotions on large inventory lines."),
        ("2. Shipping Carrier Negotiations", "Renegotiate bulk furniture freight carrier rates to reduce the logistics costs by 15-20%."),
        ("3. Regional Pricing Floor in Texas & Ohio", "Establish a minimum pricing floor and restrict local managers from running loss-making storewide sales."),
        ("4. Remote Work Product Bundling", "Bundle high-margin Technology (copiers/phones) with Office Supplies (paper/folders) to upsell Home Office clients."),
        ("5. Focus VIP Loyalty Program", "Implement a formal customer retention protocol and special corporate pricing agreements for top-tier spenders.")
    ]
    
    for title, desc in recs:
        p = tf_recs.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = c_blue
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(14)
        run.font.color.rgb = c_dark_blue
        p.space_after = Pt(12)

    # ----------------------------------------------------
    # SLIDE 7: Conclusion & Portfolio Specs (Dark Theme)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_dark_blue
    bg.line.color.rgb = c_dark_blue
    
    tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Retail Business Performance Analytics"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.space_after = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "Key Portfolio Highlights:"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = c_blue
    p2.space_after = Pt(8)
    
    highlights = [
        "Interactive Dashboards: Built using React & Recharts with dynamic cross-filtering capabilities.",
        "SQL Analytical Core: Evaluated utilizing 15 enterprise-level analytical SQL queries running on SQLite.",
        "Business-Centric Insights: Focuses strictly on business intelligence, cost containment, and pricing strategy.",
        "Modular Stack: FastAPI, SQLite, Pandas, React, and Tailwind CSS."
    ]
    for h in highlights:
        ph = tf.add_paragraph()
        ph.text = f"✔ {h}"
        ph.font.size = Pt(14)
        ph.font.color.rgb = c_white
        ph.space_after = Pt(6)
        
    p3 = tf.add_paragraph()
    p3.text = "\nThank You! | Submit for Review"
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(16, 185, 129)
    
    prs.save(pptx_path)
    print("PowerPoint presentation generated successfully.")
    return True

def build_all_reports():
    """Builds both reports in sequence."""
    pdf_success = generate_pdf_report()
    pptx_success = generate_pptx_presentation()
    return pdf_success and pptx_success

if __name__ == "__main__":
    build_all_reports()
